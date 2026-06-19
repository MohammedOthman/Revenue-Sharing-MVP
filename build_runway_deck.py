#!/usr/bin/env python3
"""
Build the Partner Revenue OS Pre-Seed Runway deck (.pptx).
All figures computed from founder-provided cost data (Pillars 1-3) + repo inputs.
No invented numbers: revenue base = 0; upside is clearly-flagged illustrative.
"""
import os
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.ticker import FuncFormatter

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

REPO = "/home/user/Revenue-Sharing-MVP"
ASSETS = os.path.join(REPO, "runway_assets")
os.makedirs(ASSETS, exist_ok=True)
OUT = os.path.join(REPO, "Partner_Revenue_OS_Pre_Seed_Runway.pptx")

# ----------------------------- MODEL --------------------------------------
START_CASH = 0
RAISE = 2_000_000
FLOOR = 150_000

P1 = 158_083.75                       # Pillar 1, one-time, Month 1
P2_payroll = 715_000 / 12             # founders
P2_gosi = 58_506.25 / 12              # GOSI & levies
P2_cloud = 55_800 / 12                # cloud
P2_monthly = P2_payroll + P2_gosi + P2_cloud
P3_monthly = 84_200.0                 # execution team, Months 2-12

MONTHS = list(range(1, 13))
cash_out = []
for m in MONTHS:
    out = P2_monthly
    if m == 1:
        out += P1
    if m >= 2:
        out += P3_monthly
    cash_out.append(out)

def run_monthly(inflows, revenue=None):
    revenue = revenue or {}
    bal, opening, ending, cin = START_CASH, [], [], []
    for m in MONTHS:
        opening.append(bal)
        inflow = inflows.get(m, 0) + revenue.get(m, 0)
        cin.append(inflow)
        bal = bal + inflow - cash_out[m - 1]
        ending.append(bal)
    return opening, cin, ending

# Scenarios
base_in = {1: RAISE}
deal_cash = 132_000          # 110k ARR annual-prepay + 22k implementation (illustrative)
deal_months = [6, 8, 10, 12]
upside_rev = {m: deal_cash for m in deal_months}

op_b, cin_b, end_base = run_monthly(base_in)
op_u, cin_u, end_up = run_monthly(base_in, upside_rev)
_, _, end_tr_on = run_monthly({1: 1_000_000, 6: 1_000_000})
_, _, end_tr_late = run_monthly({1: 1_000_000, 7: 1_000_000})

avg_burn = sum(cash_out) / 12
committed = P1 + 829_306.25 + 926_200.0
unallocated = RAISE - committed
runway_to_zero = 12 + end_base[-1] / cash_out[-1] if end_base[-1] > 0 else None

def breach_month(series):
    for i, v in enumerate(series):
        if v < FLOOR:
            return i + 1
    return None

bm_base = breach_month(end_base)
bm_up = breach_month(end_up)
bm_tr_late = breach_month(end_tr_late)

# Weekly 13-week (lump sum base): setup W1, payroll month-ends W4/W9/W13
WEEKS = list(range(1, 14))
wk_in = {1: RAISE}
wk_out = {1: P1, 4: P2_monthly, 9: P2_monthly + P3_monthly, 13: P2_monthly + P3_monthly}
def run_weekly(inflows, outs):
    bal, opening, ending, cin, cout = 0, [], [], [], []
    for w in WEEKS:
        opening.append(bal)
        i = inflows.get(w, 0); o = outs.get(w, 0)
        cin.append(i); cout.append(o)
        bal = bal + i - o
        ending.append(bal)
    return opening, cin, cout, ending
w_op, w_cin, w_cout, w_end = run_weekly(wk_in, wk_out)
# small-first-tranche stress (illustrative SAR 200k first)
_, _, _, w_end_small = run_weekly({1: 200_000}, wk_out)
min_first_tranche = P1 + P2_monthly + FLOOR  # to clear month-1 above floor

# ----------------------------- STYLE --------------------------------------
NAVY = RGBColor(0x0F, 0x2A, 0x4A)
BLUE = RGBColor(0x1C, 0x6E, 0xA4)
TEAL = RGBColor(0x18, 0x9A, 0xB4)
GREEN = RGBColor(0x1E, 0x8E, 0x3E)
RED = RGBColor(0xC0, 0x39, 0x2B)
AMBER = RGBColor(0xB8, 0x86, 0x0B)
GREY = RGBColor(0x5A, 0x63, 0x6E)
LIGHT = RGBColor(0xEE, 0xF2, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x22, 0x26, 0x2A)

def sar(x, dec=0):
    return f"SAR {x:,.{dec}f}"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def slide():
    return prs.slides.add_slide(BLANK)

def box(s, l, t, w, h):
    tb = s.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    return tb

def set_text(tf, text, size=18, bold=False, color=BLACK, align=PP_ALIGN.LEFT, italic=False):
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = "Calibri"
    return p

def rect(s, l, t, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def header(s, title, kicker=None):
    rect(s, 0, 0, SW, Inches(1.15), NAVY)
    rect(s, 0, Inches(1.15), SW, Inches(0.06), TEAL)
    tb = box(s, Inches(0.5), Inches(0.18), Inches(12.4), Inches(0.95))
    tf = tb.text_frame
    if kicker:
        set_text(tf, kicker, size=12, bold=True, color=TEAL)
        p = tf.add_paragraph(); r = p.add_run(); r.text = title
        r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Calibri"
    else:
        set_text(tf, title, size=28, bold=True, color=WHITE)

def footer(s, n):
    tb = box(s, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.35))
    p = set_text(tb.text_frame,
                 "Partner Revenue OS · Pre-Seed Runway · SAR · figures from founder cost data (Pillars 1-3); revenue base=0, upside illustrative",
                 size=8, color=GREY)
    tb2 = box(s, Inches(12.4), Inches(7.05), Inches(0.7), Inches(0.35))
    set_text(tb2.text_frame, str(n), size=9, color=GREY, align=PP_ALIGN.RIGHT)

def bullets(s, items, left=Inches(0.6), top=Inches(1.5), width=Inches(12.1), height=Inches(5.3), size=15):
    tb = box(s, left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for it in items:
        if isinstance(it, tuple):
            text, lvl, *rest = it
            color = rest[0] if rest else BLACK
            bold = rest[1] if len(rest) > 1 else False
        else:
            text, lvl, color, bold = it, 0, BLACK, False
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = lvl
        prefix = "" if lvl == 0 else ""
        bullet = "▪ " if lvl == 0 else "– "
        r = p.add_run(); r.text = bullet + text
        r.font.size = Pt(size - lvl * 2); r.font.bold = bold
        r.font.color.rgb = color; r.font.name = "Calibri"
        p.space_after = Pt(6)
    return tb

def table(s, headers, rows, top=Inches(1.45), left=Inches(0.5), width=Inches(12.33),
          height=Inches(5.2), fs=11, col_w=None, status_col=None):
    nrows, ncols = len(rows) + 1, len(headers)
    gf = s.shapes.add_table(nrows, ncols, left, top, width, height)
    t = gf.table
    if col_w:
        for i, w in enumerate(col_w):
            t.columns[i].width = w
    # header
    for j, h in enumerate(headers):
        c = t.cell(0, j); c.text = h
        c.fill.solid(); c.fill.fore_color.rgb = NAVY
        for p in c.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            for r in p.runs:
                r.font.size = Pt(fs); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Calibri"
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = t.cell(i + 1, j); c.text = str(val)
            c.fill.solid()
            c.fill.fore_color.rgb = WHITE if i % 2 == 0 else LIGHT
            color = BLACK
            if status_col is not None and j == status_col:
                if "BELOW" in str(val) or "INSOLVENT" in str(val):
                    color = RED
                elif "TIGHT" in str(val):
                    color = AMBER
                else:
                    color = GREEN
            for p in c.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
                for r in p.runs:
                    r.font.size = Pt(fs); r.font.color.rgb = color; r.font.name = "Calibri"
                    if status_col is not None and j == status_col:
                        r.font.bold = True
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.margin_left = Inches(0.05); c.margin_right = Inches(0.05)
            c.margin_top = Inches(0.01); c.margin_bottom = Inches(0.01)
    return gf

def notes(s, text):
    s.notes_slide.notes_text_frame.text = text

# ----------------------------- CHARTS -------------------------------------
def money_fmt(x, _):
    if abs(x) >= 1e6:
        return f"{x/1e6:.1f}M"
    return f"{x/1e3:.0f}k"

def line_chart(path, xs, series, xlabel, title, floor=True, xticklabels=None):
    fig, ax = plt.subplots(figsize=(8.6, 4.2), dpi=150)
    colors = ["#1C6EA4", "#1E8E3E", "#C0392B", "#B8860B"]
    for i, (label, ys) in enumerate(series):
        ax.plot(xs, ys, marker="o", linewidth=2.4, label=label, color=colors[i % len(colors)])
    if floor:
        ax.axhline(FLOOR, color="#C0392B", linestyle="--", linewidth=1.6, label=f"Floor SAR {FLOOR:,.0f}")
        ax.axhline(0, color="#888", linewidth=1)
    ax.set_xlabel(xlabel, fontsize=10)
    ax.set_ylabel("Ending cash (SAR)", fontsize=10)
    ax.set_title(title, fontsize=12, fontweight="bold", color="#0F2A4A")
    ax.yaxis.set_major_formatter(FuncFormatter(money_fmt))
    if xticklabels:
        ax.set_xticks(xs); ax.set_xticklabels(xticklabels, fontsize=8)
    ax.grid(True, alpha=0.25)
    ax.legend(fontsize=8, loc="best")
    ax.spines["top"].set_visible(False); ax.spines["right"].set_visible(False)
    fig.tight_layout(); fig.savefig(path, bbox_inches="tight"); plt.close(fig)

def pie_chart(path):
    fig, ax = plt.subplots(figsize=(5.6, 4.4), dpi=150)
    vals = [P1, 829_306.25, 926_200.0, unallocated]
    labels = [f"Setup\n{P1/1e3:.0f}k (7.9%)",
              f"Founders+Cloud\n829k (41.5%)",
              f"Execution team\n926k (46.3%)",
              f"Unallocated\n{unallocated/1e3:.0f}k (4.3%)"]
    colors = ["#1C6EA4", "#189AB4", "#1E8E3E", "#C9CDD2"]
    ax.pie(vals, labels=labels, colors=colors, autopct="", startangle=90,
           textprops={"fontsize": 9}, wedgeprops={"edgecolor": "white", "linewidth": 1.5})
    ax.set_title("How the SAR 2.0M is allocated", fontsize=12, fontweight="bold", color="#0F2A4A")
    fig.tight_layout(); fig.savefig(path, bbox_inches="tight"); plt.close(fig)

c_monthly = os.path.join(ASSETS, "monthly.png")
c_tranche = os.path.join(ASSETS, "tranche.png")
c_weekly = os.path.join(ASSETS, "weekly.png")
c_pie = os.path.join(ASSETS, "pie.png")
mlabels = [f"M{m}" for m in MONTHS]
line_chart(c_monthly, MONTHS, [("Base (no revenue)", end_base), ("Illustrative upside", end_up)],
           "Month", "12-month ending cash: base vs illustrative upside", xticklabels=mlabels)
line_chart(c_tranche, MONTHS,
           [("Lump sum (2.0M @ M1)", end_base),
            ("Tranches on-time (1.0M M1 + 1.0M M6)", end_tr_on),
            ("Tranche 2 slips to M7", end_tr_late)],
           "Month", "Inflow timing stress test", xticklabels=mlabels)
line_chart(c_weekly, WEEKS, [("Lump sum base", w_end), ("Small first tranche (200k)", w_end_small)],
           "Week", "13-week weekly cash (lump sum vs too-small first tranche)",
           xticklabels=[f"W{w}" for w in WEEKS])
pie_chart(c_pie)

# ----------------------------- SLIDES -------------------------------------
def status_for(v):
    if v < 0: return "INSOLVENT"
    if v < FLOOR: return "BELOW FLOOR"
    if v < FLOOR * 1.5: return "TIGHT"
    return "OK"

# 1. TITLE
s = slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(3.05), SW, Inches(0.08), TEAL)
tb = box(s, Inches(0.8), Inches(1.7), Inches(11.7), Inches(1.4))
set_text(tb.text_frame, "Partner Revenue OS", size=40, bold=True, color=WHITE)
tb2 = box(s, Inches(0.8), Inches(3.25), Inches(11.7), Inches(1.6))
tf = tb2.text_frame
set_text(tf, "Pre-Seed Runway — 13-Week Weekly + 12-Month Cash-Flow Model", size=24, bold=True, color=TEAL)
p = tf.add_paragraph(); r = p.add_run()
r.text = "Zero-revenue base case + illustrative upside · SAR · KSA/GCC-calibrated"
r.font.size = Pt(16); r.font.color.rgb = WHITE; r.font.name = "Calibri"
tb3 = box(s, Inches(0.8), Inches(5.7), Inches(11.7), Inches(1.2))
tf = tb3.text_frame
set_text(tf, "June 2026  ·  Built only from founder-provided cost data (Pillars 1-3) and the repo's regulatory reference.", size=12, color=RGBColor(0xC8,0xD2,0xDC))
p = tf.add_paragraph(); r = p.add_run()
r.text = "No invented figures: every number is founder input (real) or a clearly-flagged assumption."
r.font.size = Pt(12); r.font.italic = True; r.font.color.rgb = RGBColor(0xC8,0xD2,0xDC); r.font.name="Calibri"
notes(s, "This deck answers: once the SAR 2.0M lands, how long does it last, which week/month gets tight, "
         "and what must change. Base case assumes zero revenue (the honest survival view). The upside is "
         "illustrative, built from the repo's pricing ranges, and clearly flagged.")

# 2. EXECUTIVE SUMMARY
s = slide(); header(s, "Executive summary", "THE BOTTOM LINE")
bullets(s, [
    (f"Starting cash is SAR 0. The clock starts when the first capital lands. Raise = SAR 2.0M (timing not yet decided).", 0, NAVY, True),
    (f"The SAR 2.0M is ~95.7% pre-committed: Setup + 4 founders + 6 execution hires + cloud = SAR {committed:,.0f}. Only SAR {unallocated:,.0f} (4.3%) is uncommitted.", 0),
    (f"Average burn ≈ SAR {avg_burn:,.0f}/month. Month 1 is front-loaded at SAR {cash_out[0]:,.0f} (setup + first payroll).", 0),
    (f"BASE CASE (zero revenue, lump sum): cash lasts the year but ends at SAR {end_base[-1]:,.0f} — BELOW the SAR 150k floor, breached around Month {bm_base}. Runway-to-zero ≈ {runway_to_zero:.1f} months.", 0, RED, True),
    (f"This is 'default-dead' on cost alone: the raise funds the team for ~a year with no cushion, no tax line, and no marketing budget.", 1, RED),
    (f"ILLUSTRATIVE UPSIDE (3-5 mid-market logos in H2): ends the year at ~SAR {end_up[-1]:,.0f}, stays above the floor — revenue is the difference between survival and stall.", 0, GREEN, True),
    (f"If funded in tranches, the FIRST tranche must be ≥ ~SAR {min_first_tranche:,.0f} or Month 1 breaches the floor; tranche 2 must land by ~Month 6.", 0, AMBER, True),
    ("Recommendation: front-load the raise (lump sum or a large first tranche), itemize taxes + GTM, and treat H2 revenue as the gating milestone for survival and the seed round.", 0, NAVY, True),
], size=14)
footer(s, 2)
notes(s, "One-screen verdict for a board. The headline tension: the cost plan is fully specified and consumes "
         "almost the entire raise; survival past ~month 11-12 depends on revenue and/or more capital.")

# 3. METHOD & RULES
s = slide(); header(s, "How to read this model", "METHOD & DISCIPLINE")
bullets(s, [
    ("This is a CASH model — it tracks money in and money out of the bank, week by week and month by month.", 0, NAVY, True),
    ("Opening cash  +  cash in  −  cash out  =  ending cash.  Ending cash carries to the next period. Repeat.", 1),
    ("Every figure is tagged: REAL (founder-provided cost data) or ASSUMPTION (clearly flagged, editable).", 0),
    ("No invented numbers. Where an input isn't decided yet, the model uses a stated assumption you can overwrite.", 1),
    ("Four things kept separate: bookings ≠ invoices ≠ recognised revenue ≠ CASH COLLECTED. This model = cash only.", 0),
    ("VAT is treated as working capital, never revenue. Currency = SAR.", 1),
    ("Sources: founder cost sheet (Pillars 1-3, the verified financial model) + the repo's KSA/GCC regulatory reference for rates and timing.", 0),
    ("Scope: all 12 months sit inside Phase 1 'Capture' (the software product) — no money-movement / settlement in this window.", 0),
], size=14)
footer(s, 3)
notes(s, "Sets expectations: cash not accounting profit; real vs assumed; the no-hallucination discipline from the repo's own model prompt.")

# 4. INPUTS / ASSUMPTIONS REGISTER
s = slide(); header(s, "Inputs & assumptions register", "WHAT GOES IN")
table(s,
    ["Input", "Value", "Tier", "Note"],
    [
        ["Starting cash", "SAR 0", "REAL", "Clock starts at first capital injection"],
        ["Raise", "SAR 2,000,000", "REAL", "Instrument (shares/SAFE) & timing undecided"],
        ["Inflow timing", "Lump sum @ M1", "ASSUMPTION", "Base case; tranche stress shown separately"],
        ["Cash floor", "SAR 150,000", "REAL", "Do-not-breach line"],
        ["Team", "4 founders (M1) + 6 hires (M2)", "REAL", "10 people by Month 2; mixed Saudi/expat"],
        ["Pillar 1 (setup)", "SAR 158,083.75", "REAL", "One-time, Month 1"],
        ["Pillar 2 (founders+cloud)", "SAR 829,306.25 / yr", "REAL", "Payroll 715k + GOSI 58.5k + cloud 55.8k"],
        ["Pillar 3 (execution)", "SAR 926,200 / 11 mo", "REAL", "SAR 84,200/mo from Month 2"],
        ["Revenue (base)", "SAR 0", "ASSUMPTION", "Honest survival view; nothing signed"],
        ["Revenue (upside)", "~SAR 132k/deal, H2", "ILLUSTRATIVE", "110k ARR prepay + 22k impl; repo pricing ranges"],
        ["Taxes (Zakat/VAT/WHT)", "Not yet modelled", "GAP", "Need ownership split, VAT status, foreign vendors"],
    ],
    fs=11, col_w=[Inches(3.0), Inches(2.7), Inches(1.7), Inches(4.9)], status_col=2)
footer(s, 4)
notes(s, "The single editable control panel. REAL = founder data. ASSUMPTION/ILLUSTRATIVE/GAP are what to refine.")

# 5. COST STRUCTURE (pie + text)
s = slide(); header(s, "The cost structure: where the SAR 2.0M goes", "PILLARS 1-3")
s.shapes.add_picture(c_pie, Inches(0.4), Inches(1.5), height=Inches(4.9))
bullets(s, [
    ("Setup (Pillar 1): SAR 158,084 one-time — CR/MISA, visas & Iqamas, portals, office, accounting+insurance.", 0),
    ("Founders + cloud (Pillar 2): SAR 829,306/yr — 4 founders, GOSI/levies, Riyadh cloud.", 0),
    ("Execution team (Pillar 3): SAR 926,200 — 6 hires from Month 2 (architect, DevOps, integration, sales, CSM, support).", 0),
    ("Committed = SAR 1,913,590 = 95.7% of the raise.", 0, NAVY, True),
    ("Unallocated = SAR 86,410 (4.3%) — the only cushion.", 0, AMBER, True),
    ("Not yet in any pillar: marketing/GTM spend, taxes, travel, contingency.", 0, RED, True),
], left=Inches(6.5), top=Inches(1.7), width=Inches(6.4), size=13)
footer(s, 5)
notes(s, "Visual of the commitment problem: nearly the whole raise is people + setup; no discretionary room.")

# 6. MONTHLY RUN-RATE BUILD
s = slide(); header(s, "Monthly run-rate: how fast cash leaves", "THE BURN")
table(s,
    ["Cost line", "SAR / month", "When"],
    [
        ["Founder payroll (4)", f"{P2_payroll:,.2f}", "Every month, M1-12"],
        ["GOSI & statutory levies", f"{P2_gosi:,.2f}", "Every month, M1-12"],
        ["Cloud infrastructure", f"{P2_cloud:,.2f}", "Every month, M1-12"],
        ["Execution team (6)", f"{P3_monthly:,.2f}", "From Month 2 (M2-12)"],
        ["Setup (one-time)", f"{P1:,.2f}", "Month 1 only"],
        ["— Run-rate Month 1", f"{cash_out[0]:,.2f}", "Front-loaded (setup + payroll)"],
        ["— Run-rate Months 2-12", f"{cash_out[1]:,.2f}", "Steady monthly burn"],
        ["— Average burn / month", f"{avg_burn:,.2f}", "Across the 12 months"],
    ],
    fs=12.5, col_w=[Inches(4.6), Inches(3.6), Inches(4.13)])
footer(s, 6)
notes(s, "Two regimes: a heavy Month 1 (setup lands at once), then a flat ~153k/month once the team is on.")

# 7. 12-MONTH BASE TABLE
s = slide(); header(s, "12-month runway — base case (zero revenue, lump sum)", "MONTHLY")
rows = []
for i, m in enumerate(MONTHS):
    rows.append([f"M{m}", f"{cin_b[i]:,.0f}", f"{cash_out[i]:,.0f}", f"{end_base[i]:,.0f}", status_for(end_base[i])])
table(s, ["Month", "Cash in", "Cash out", "Ending cash", "Status"], rows,
      fs=11.5, col_w=[Inches(1.6), Inches(2.7), Inches(2.7), Inches(2.9), Inches(2.43)],
      top=Inches(1.4), height=Inches(5.2), status_col=4)
footer(s, 7)
notes(s, f"With the full 2.0M up front and no revenue, cash declines steadily, drops below the 150k floor in "
         f"Month {bm_base}, and ends the year at SAR {end_base[-1]:,.0f}. It never hits zero inside 12 months, "
         f"but only just (runway-to-zero ~{runway_to_zero:.1f} months) — and this ignores tax and GTM.")

# 8. MONTHLY CHART base vs upside
s = slide(); header(s, "Base vs illustrative upside", "12-MONTH ENDING CASH")
s.shapes.add_picture(c_monthly, Inches(0.5), Inches(1.45), width=Inches(8.6))
bullets(s, [
    ("Red dashed line = SAR 150k floor.", 0, GREY),
    (f"Base (no revenue) ends at SAR {end_base[-1]:,.0f} — below floor.", 0, RED, True),
    (f"Upside ends at SAR {end_up[-1]:,.0f} — above floor.", 0, GREEN, True),
    ("Upside = 4 mid-market logos landing in H2 (illustrative, from the repo's pricing ranges).", 0),
    ("Takeaway: a handful of paid logos is the difference between stalling and reaching the seed bar.", 0, NAVY, True),
], left=Inches(9.3), top=Inches(1.7), width=Inches(3.7), size=12)
footer(s, 8)
notes(s, "The visual case for prioritising revenue in H2. Upside deliberately conservative on timing (long enterprise cycles).")

# 9. 13-WEEK WEEKLY
s = slide(); header(s, "13-week weekly cash (the near-term liquidity view)", "WEEKLY")
rows = []
for i, w in enumerate(WEEKS):
    rows.append([f"W{w}", f"{w_cin[i]:,.0f}" if w_cin[i] else "—",
                 f"{w_cout[i]:,.0f}" if w_cout[i] else "—", f"{w_end[i]:,.0f}", status_for(w_end[i])])
table(s, ["Week", "Cash in", "Cash out", "Ending cash", "Status"], rows,
      fs=10.5, col_w=[Inches(1.5), Inches(2.6), Inches(2.6), Inches(2.8), Inches(2.0)],
      top=Inches(1.4), height=Inches(5.2), status_col=4)
footer(s, 9)
notes(s, "Lump-sum case: setup hits in W1, payroll at month-ends (W4/W9/W13). The 13-week window is the calm "
         "early period; the crunch (in the tranche case) comes later, around month 5-6.")

# 10. WEEKLY CHART
s = slide(); header(s, "13-week weekly cash — why the first tranche size matters", "WEEKLY")
s.shapes.add_picture(c_weekly, Inches(0.5), Inches(1.45), width=Inches(8.6))
bullets(s, [
    ("Lump sum: setup in W1 leaves ~SAR 1.84M; calm through the quarter (ends ~SAR 1.47M).", 0, GREEN, True),
    (f"A too-small first tranche (e.g. SAR 200k) breaches the floor in Week 1 (ends SAR {w_end_small[0]:,.0f}).", 0, RED, True),
    (f"First tranche must clear setup + 1 month payroll + floor ≈ SAR {min_first_tranche:,.0f}.", 0, AMBER, True),
    ("This is the near-term liquidity early-warning the weekly view exists to give.", 0, NAVY),
], left=Inches(9.3), top=Inches(1.7), width=Inches(3.7), size=12)
footer(s, 10)
notes(s, "Weekly chart visualising the first-tranche threshold and the immediate breach if it's too small.")

# 11. TRANCHE STRESS
s = slide(); header(s, "Inflow timing: the single biggest risk", "STRESS TEST")
s.shapes.add_picture(c_tranche, Inches(0.5), Inches(1.45), width=Inches(8.6))
bullets(s, [
    ("Start cash is zero, so survival depends entirely on WHEN the 2.0M arrives.", 0, NAVY, True),
    (f"Minimum first tranche ≈ SAR {min_first_tranche:,.0f} (setup + 1 month payroll + floor).", 0, AMBER, True),
    ("A too-small first tranche breaches the floor in Week 1 (see weekly chart).", 0, RED),
    ("With a 50/50 split, tranche 2 must land by ~Month 6.", 0),
    (f"If tranche 2 slips one month to M7, cash nearly hits zero (breach at Month {bm_tr_late}).", 0, RED, True),
    ("Mitigation: lump sum, or a large front-loaded first tranche with wide spacing.", 0, GREEN, True),
], left=Inches(9.3), top=Inches(1.7), width=Inches(3.7), size=12)
footer(s, 11)
notes(s, "Quantifies why front-loading matters. Tranche assumptions are illustrative (50/50) to show sensitivity.")

# 11. UPSIDE DETAIL
s = slide(); header(s, "Illustrative upside — what revenue does", "SCENARIO (FLAGGED)")
bullets(s, [
    ("ALL figures here are illustrative assumptions from the repo's pricing ranges — nothing is signed.", 0, AMBER, True),
    ("Assumed motion: the Inside Sales Executive (hired M2) + founders land mid-market logos in H2 after ramp + sales cycle.", 0),
    ("Deal economics (illustrative): SAR ~110,000 annual contract, paid up front (GCC norm) + ~SAR 22,000 implementation = ~SAR 132,000 cash per logo.", 0),
    ("Timing: one logo each in Months 6, 8, 10, 12 (≈ the repo's 3-5 paid design-partner seed milestone).", 0),
    (f"Cash impact: +SAR {sum(upside_rev.values()):,.0f} across H2 → year-end cash ~SAR {end_up[-1]:,.0f} (vs {end_base[-1]:,.0f} base).", 0, GREEN, True),
    ("Year-end ARR ≈ SAR 440,000 (4 logos) — entering seed-ready territory.", 0, GREEN),
    ("Reality check: enterprise cycles are long; treat these dates as optimistic and pressure-test pilot→paid conversion.", 0, GREY, True),
], size=13)
footer(s, 12)
notes(s, "Keeps the upside honest: clearly labelled illustrative, conservative on count and timing, tied to the repo's own seed milestone.")

# 12. GAPS & RISKS
s = slide(); header(s, "Gaps to close (ranked by impact)", "WHAT'S MISSING")
table(s,
    ["#", "Gap", "Why it matters", "Impact"],
    [
        ["1", "Inflow timing of the 2.0M", "Zero start cash → survival hinges on when money lands", "HIGH"],
        ["2", "Revenue (none decided)", "Sets the real cash-out month; flips default-dead → alive", "HIGH"],
        ["3", "Taxes & timing (Zakat/VAT/WHT)", "Real cash outflows; need ownership %, VAT status, vendors", "MED-HIGH"],
        ["4", "Pillar 3 burden inclusion", "Do role budgets include GOSI/levies, or are they on top?", "MEDIUM"],
        ["5", "Missing GTM / contingency budget", "No marketing/travel line; 86k is the only cushion", "MEDIUM"],
        ["6", "Exact pay dates / week timing", "Refines the weekly view (currently month-end default)", "LOW"],
        ["7", "715k vs 720k salary; Saudi/expat split", "Minor reconciliations", "LOW"],
    ],
    fs=11.5, col_w=[Inches(0.6), Inches(3.5), Inches(6.4), Inches(1.83)], status_col=3)
footer(s, 13)
notes(s, "The honest 'to make this board-defensible' list. Items 1-3 move the conclusion; the rest refine accuracy.")

# 13. RECOMMENDATIONS
s = slide(); header(s, "Recommendations & decision gates", "WHAT TO DO")
bullets(s, [
    ("CAPITAL: take the 2.0M as a lump sum or a large front-loaded first tranche (≥ ~SAR 377k); avoid thin/late tranches.", 0, NAVY, True),
    ("RUNWAY: plan to the floor, not to zero — at current burn the 150k floor hits around Month 11-12 with no revenue.", 0),
    ("REVENUE: make H2 paid logos the gating milestone; even 3-5 mid-market deals change survival and unlock the seed round.", 0, GREEN, True),
    ("BUDGET: itemize the missing lines (taxes, GTM, contingency) and decide if GTM comes from this raise or the next.", 0),
    ("RAISE PLAN: this round funds ~12 months with no cushion → start the seed raise by ~Month 6-7 (MENA processes run 6-9 months).", 0, AMBER, True),
    ("VERIFY: confirm ownership split (Zakat vs CIT), VAT registration, and WHT on foreign cloud/SaaS before filing.", 0),
    ("GOVERNANCE: review cash weekly for the first 6 months; re-run this model monthly against actuals.", 0),
], size=14)
footer(s, 14)
notes(s, "Actionable close: capital structure, the revenue milestone, the missing budget, and when to start the next raise.")

# 14. APPENDIX — upside table
s = slide(); header(s, "Appendix — monthly detail (illustrative upside)", "APPENDIX")
rows = []
for i, m in enumerate(MONTHS):
    rev = upside_rev.get(m, 0)
    rows.append([f"M{m}", f"{rev:,.0f}" if rev else "—", f"{cin_u[i]:,.0f}",
                 f"{cash_out[i]:,.0f}", f"{end_up[i]:,.0f}", status_for(end_up[i])])
table(s, ["Month", "Revenue cash", "Total in", "Cash out", "Ending cash", "Status"], rows,
      fs=11, col_w=[Inches(1.5), Inches(2.4), Inches(2.4), Inches(2.4), Inches(2.4), Inches(1.6)],
      top=Inches(1.4), height=Inches(5.2), status_col=5)
footer(s, 15)
notes(s, "Full month-by-month for the illustrative upside, for reference.")

prs.save(OUT)

# verification printout
print("=== KEY FIGURES ===")
print("committed", f"{committed:,.2f}", "unallocated", f"{unallocated:,.2f}")
print("M1 out", f"{cash_out[0]:,.2f}", "M2-12 out", f"{cash_out[1]:,.2f}", "avg burn", f"{avg_burn:,.2f}")
print("base ending M12", f"{end_base[-1]:,.2f}", "breach month", bm_base, "runway", f"{runway_to_zero:.2f}")
print("upside ending M12", f"{end_up[-1]:,.2f}", "breach", bm_up)
print("tranche late breach month", bm_tr_late)
print("min first tranche", f"{min_first_tranche:,.2f}")
print("weekly end W13", f"{w_end[-1]:,.2f}", "weekly small-first W1", f"{w_end_small[0]:,.2f}")
print("saved", OUT, os.path.getsize(OUT), "bytes")
